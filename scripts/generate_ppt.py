#!/usr/bin/env python3
"""
Narae TMS v2.0 소개 PPT 생성 스크립트
python-pptx를 사용하여 18슬라이드 와이드스크린(16:9) PPT를 생성합니다.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ──────────────────────────────────────────────
# Color Theme
# ──────────────────────────────────────────────
PRIMARY_DARK = RGBColor(0x0F, 0x17, 0x2A)
PRIMARY_BLUE = RGBColor(0x25, 0x63, 0xEB)
PRIMARY_PURPLE = RGBColor(0x93, 0x33, 0xEA)
PRIMARY_CYAN = RGBColor(0x06, 0xB6, 0xD4)
SUCCESS_GREEN = RGBColor(0x22, 0xC5, 0x5E)
BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_PRIMARY = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE_BG = RGBColor(0xEF, 0xF6, 0xFF)
LIGHT_PURPLE_BG = RGBColor(0xF5, 0xF3, 0xFF)
LIGHT_GREEN_BG = RGBColor(0xF0, 0xFD, 0xF4)
LIGHT_CYAN_BG = RGBColor(0xEC, 0xFE, 0xFF)
LIGHT_ORANGE_BG = RGBColor(0xFF, 0xF7, 0xED)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
RED_SOFT = RGBColor(0xEF, 0x44, 0x44)

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Common dimensions
MARGIN_LEFT = Inches(0.8)
MARGIN_RIGHT = Inches(0.8)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
TITLE_BAR_HEIGHT = Inches(1.0)
FOOTER_HEIGHT = Inches(0.4)

COMPANY_NAME = "주식회사 나래정보기술"
PRODUCT_NAME = "Narae TMS v2.0"
TAGLINE = "Enterprise-Grade Oracle SQL Tuning Management System"


def hex_to_rgb(hex_str):
    """Convert hex color string to RGBColor."""
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ──────────────────────────────────────────────
# Helper Functions
# ──────────────────────────────────────────────

def set_slide_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
        shape.line.fill.solid()
    # Minimal corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_rect(slide, left, top, width, height, fill_color=None):
    """Add a simple rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=14,
                font_color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="맑은 고딕"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=13,
                       font_color=TEXT_PRIMARY, line_spacing=1.3, font_name="맑은 고딕",
                       bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    font_color=TEXT_PRIMARY, bullet_char="•", font_name="맑은 고딕",
                    line_spacing=1.4):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return txBox


def add_title_bar(slide, title, subtitle=None, accent_color=PRIMARY_BLUE):
    """Add a common title bar at the top of the slide."""
    # Accent line
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), accent_color)
    # Title background
    add_rect(slide, Inches(0), Inches(0.06), SLIDE_WIDTH, Inches(0.94), WHITE)
    # Title text
    add_textbox(slide, MARGIN_LEFT, Inches(0.18), Inches(10), Inches(0.5),
                title, font_size=26, font_color=PRIMARY_DARK, bold=True)
    if subtitle:
        add_textbox(slide, MARGIN_LEFT, Inches(0.58), Inches(10), Inches(0.35),
                    subtitle, font_size=13, font_color=TEXT_MUTED)
    # Divider line
    add_rect(slide, Inches(0), Inches(1.0), SLIDE_WIDTH, Inches(0.015), CARD_BORDER)


def add_footer(slide, page_num, total=18):
    """Add footer with page number and company name."""
    y = SLIDE_HEIGHT - FOOTER_HEIGHT
    add_rect(slide, Inches(0), y, SLIDE_WIDTH, FOOTER_HEIGHT, WHITE)
    add_rect(slide, Inches(0), y, SLIDE_WIDTH, Inches(0.01), CARD_BORDER)
    add_textbox(slide, MARGIN_LEFT, y + Inches(0.05), Inches(4), Inches(0.3),
                f"{PRODUCT_NAME}  |  {COMPANY_NAME}", font_size=8, font_color=TEXT_MUTED)
    add_textbox(slide, SLIDE_WIDTH - Inches(1.8), y + Inches(0.05), Inches(1.0), Inches(0.3),
                f"{page_num} / {total}", font_size=8, font_color=TEXT_MUTED,
                alignment=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title="", items=None,
             accent_color=PRIMARY_BLUE, bg_color=WHITE, title_size=15,
             item_size=12, icon_text=None):
    """Add a card with accent bar, title, and bullet items."""
    # Card background
    card = add_shape(slide, left, top, width, height, fill_color=bg_color,
                     border_color=CARD_BORDER, border_width=Pt(1))
    # Accent bar (left side)
    add_rect(slide, left + Inches(0.02), top + Inches(0.15), Inches(0.06),
             height - Inches(0.3), accent_color)
    # Icon/emoji
    if icon_text:
        add_textbox(slide, left + Inches(0.25), top + Inches(0.12), Inches(0.4), Inches(0.4),
                    icon_text, font_size=18, font_color=accent_color, bold=True)
        title_left = left + Inches(0.65)
    else:
        title_left = left + Inches(0.25)
    # Title
    add_textbox(slide, title_left, top + Inches(0.15), width - Inches(0.8), Inches(0.35),
                title, font_size=title_size, font_color=PRIMARY_DARK, bold=True)
    # Items
    if items:
        item_top = top + Inches(0.55)
        for item in items:
            add_textbox(slide, left + Inches(0.3), item_top, width - Inches(0.6), Inches(0.28),
                        f"▸  {item}", font_size=item_size, font_color=TEXT_PRIMARY)
            item_top += Inches(0.28)
    return card


def add_metric_box(slide, left, top, width, height, value, label,
                   accent_color=PRIMARY_BLUE):
    """Add a metric display box."""
    box = add_shape(slide, left, top, width, height, fill_color=WHITE,
                    border_color=CARD_BORDER, border_width=Pt(1))
    add_textbox(slide, left, top + Inches(0.1), width, Inches(0.45),
                value, font_size=28, font_color=accent_color, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.55), width, Inches(0.3),
                label, font_size=10, font_color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)
    return box


def make_slide(prs):
    """Create a new blank slide."""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    return slide


# ──────────────────────────────────────────────
# Slide Builders
# ──────────────────────────────────────────────

def slide_01_cover(prs):
    """표지 슬라이드"""
    slide = make_slide(prs)
    # Dark gradient background
    set_slide_bg(slide, PRIMARY_DARK)

    # Decorative accent shapes
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_HEIGHT, PRIMARY_BLUE)
    add_shape(slide, Inches(9.5), Inches(-1), Inches(5), Inches(5),
              fill_color=RGBColor(0x1E, 0x29, 0x3B))
    add_shape(slide, Inches(10), Inches(4), Inches(4.5), Inches(4.5),
              fill_color=RGBColor(0x1E, 0x29, 0x3B))

    # Product name
    add_textbox(slide, Inches(1.2), Inches(1.8), Inches(8), Inches(0.6),
                "Narae TMS v2.0", font_size=48, font_color=WHITE, bold=True)

    # Tagline in Korean
    add_textbox(slide, Inches(1.2), Inches(2.65), Inches(9), Inches(0.5),
                "Oracle 데이터베이스 SQL 튜닝 관리 시스템",
                font_size=24, font_color=PRIMARY_CYAN, bold=True)

    # English subtitle
    add_textbox(slide, Inches(1.2), Inches(3.25), Inches(9), Inches(0.5),
                "Enterprise-Grade SQL Tuning Management System",
                font_size=16, font_color=TEXT_MUTED)

    # Separator line
    add_rect(slide, Inches(1.2), Inches(4.0), Inches(2.5), Inches(0.04), PRIMARY_BLUE)

    # Key metrics
    metrics = [
        ("62+", "페이지"),
        ("22", "기능 카테고리"),
        ("AI/LLM", "통합"),
        ("24/7", "모니터링"),
    ]
    for i, (val, label) in enumerate(metrics):
        x = Inches(1.2 + i * 2.2)
        add_textbox(slide, x, Inches(4.35), Inches(1.8), Inches(0.45),
                    val, font_size=28, font_color=PRIMARY_BLUE, bold=True)
        add_textbox(slide, x, Inches(4.8), Inches(1.8), Inches(0.3),
                    label, font_size=12, font_color=TEXT_MUTED)

    # Company info
    add_textbox(slide, Inches(1.2), Inches(5.8), Inches(8), Inches(0.4),
                COMPANY_NAME, font_size=16, font_color=WHITE, bold=True)
    add_textbox(slide, Inches(1.2), Inches(6.2), Inches(8), Inches(0.3),
                "www.narae-it.co.kr  |  2026", font_size=12, font_color=TEXT_MUTED)


def slide_02_problem(prs):
    """문제 정의 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "왜 SQL 튜닝 관리 시스템이 필요한가?", "현재 기업이 직면한 5가지 핵심 과제", RED_SOFT)
    add_footer(slide, 2)

    problems = [
        ("01", "수동 모니터링의 한계",
         "수십~수백 개의 SQL을 수동으로 추적하는 것은 비효율적이며, 성능 저하를 사전에 감지하기 어렵습니다.",
         RED_SOFT),
        ("02", "튜닝 이력 관리 부재",
         "튜닝 전/후 성능 비교, 변경 이력, 담당자 추적이 체계적으로 관리되지 않아 반복적인 문제 발생.",
         ORANGE),
        ("03", "전문가 의존성",
         "SQL 튜닝은 고도의 전문 지식이 필요하여, 특정 인력에 대한 의존도가 높고 지식 공유가 어렵습니다.",
         PRIMARY_PURPLE),
        ("04", "실시간 대응 불가",
         "성능 이슈 발생 시 원인 파악까지 수 시간이 소요되며, 비즈니스에 직접적인 영향을 미칩니다.",
         PRIMARY_BLUE),
        ("05", "보안 및 감사 미흡",
         "DB 접속 정보 관리, 작업 이력 추적, 권한 제어가 미흡하여 보안 리스크가 존재합니다.",
         SUCCESS_GREEN),
    ]

    y_start = Inches(1.3)
    card_height = Inches(1.05)
    gap = Inches(0.12)

    for i, (num, title, desc, color) in enumerate(problems):
        y = y_start + i * (card_height + gap)
        # Card
        add_shape(slide, MARGIN_LEFT, y, CONTENT_WIDTH, card_height,
                  fill_color=WHITE, border_color=CARD_BORDER, border_width=Pt(1))
        # Number badge
        add_shape(slide, MARGIN_LEFT + Inches(0.2), y + Inches(0.2),
                  Inches(0.6), Inches(0.6), fill_color=color)
        add_textbox(slide, MARGIN_LEFT + Inches(0.2), y + Inches(0.22),
                    Inches(0.6), Inches(0.55), num, font_size=20,
                    font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Title
        add_textbox(slide, MARGIN_LEFT + Inches(1.1), y + Inches(0.12),
                    Inches(4), Inches(0.35), title, font_size=16,
                    font_color=PRIMARY_DARK, bold=True)
        # Description
        add_textbox(slide, MARGIN_LEFT + Inches(1.1), y + Inches(0.48),
                    CONTENT_WIDTH - Inches(1.5), Inches(0.5), desc,
                    font_size=12, font_color=TEXT_MUTED)


def slide_03_solution(prs):
    """솔루션 개요 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "Narae TMS v2.0 - 6대 핵심 가치",
                  "엔터프라이즈급 Oracle SQL 성능 최적화 통합 플랫폼")
    add_footer(slide, 3)

    values = [
        ("실시간 모니터링", "Top SQL, Wait Events,\nSession, Lock 실시간 추적", PRIMARY_BLUE, LIGHT_BLUE_BG),
        ("AI/LLM 통합", "온프레미스 LLM 기반\n자동 SQL 튜닝 가이드", PRIMARY_PURPLE, LIGHT_PURPLE_BG),
        ("튜닝 워크플로우", "칸반 보드 기반 5단계\n체계적 튜닝 프로세스", PRIMARY_CYAN, LIGHT_CYAN_BG),
        ("성능 히스토리", "A-F 등급 시스템으로\n일별/시간대별 추이 분석", SUCCESS_GREEN, LIGHT_GREEN_BG),
        ("AWR/ADDM 분석", "Enterprise & Standard\nEdition 모두 지원", ORANGE, LIGHT_ORANGE_BG),
        ("보안 & 접근 제어", "JWT, bcrypt, RBAC,\nAES-256 암호화", RED_SOFT, RGBColor(0xFF, 0xF1, 0xF2)),
    ]

    cols = 3
    card_w = Inches(3.6)
    card_h = Inches(2.4)
    gap_x = Inches(0.4)
    gap_y = Inches(0.3)
    start_x = (SLIDE_WIDTH - (cols * card_w + (cols - 1) * gap_x)) // 2
    start_y = Inches(1.4)

    for i, (title, desc, color, bg) in enumerate(values):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        add_shape(slide, x, y, card_w, card_h, fill_color=bg,
                  border_color=CARD_BORDER, border_width=Pt(1))
        # Color accent bar top
        add_rect(slide, x + Inches(0.1), y + Inches(0.08), Inches(0.8), Inches(0.06), color)
        # Title
        add_textbox(slide, x + Inches(0.2), y + Inches(0.3), card_w - Inches(0.4), Inches(0.4),
                    title, font_size=16, font_color=PRIMARY_DARK, bold=True,
                    alignment=PP_ALIGN.LEFT)
        # Description
        add_multiline_text(slide, x + Inches(0.2), y + Inches(0.85),
                           card_w - Inches(0.4), Inches(1.2),
                           desc.split('\n'), font_size=12, font_color=TEXT_MUTED,
                           line_spacing=1.5)


def slide_04_architecture(prs):
    """시스템 아키텍처 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "시스템 아키텍처", "4계층 모던 웹 아키텍처 (Client → App → AI → Data)")
    add_footer(slide, 4)

    layers = [
        ("Client Layer", "Next.js 15 + React 19 + Tailwind CSS + Shadcn UI",
         ["App Router (SSR/CSR)", "React Query 상태 관리", "Zustand 전역 상태",
          "Framer Motion 애니메이션", "다크모드 지원"],
         PRIMARY_BLUE, LIGHT_BLUE_BG),
        ("Application Layer", "Next.js API Routes + NextAuth + Drizzle ORM",
         ["RESTful API Endpoints", "JWT 인증/인가", "Drizzle ORM 타입 안전성",
          "React Hook Form + Zod", "서버 컴포넌트"],
         PRIMARY_CYAN, LIGHT_CYAN_BG),
        ("AI/LLM Layer", "Ollama / vLLM (온프레미스)",
         ["SQL 튜닝 가이드 생성", "성능 패턴 분석", "자연어 SQL 검색",
          "스트리밍 응답", "Clean Architecture"],
         PRIMARY_PURPLE, LIGHT_PURPLE_BG),
        ("Data Layer", "PostgreSQL 17 + Oracle Database",
         ["PostgreSQL: 앱 메타데이터", "Oracle: 대상 DB 모니터링",
          "AES-256 연결 암호화", "Drizzle 마이그레이션", "감사 로그"],
         SUCCESS_GREEN, LIGHT_GREEN_BG),
    ]

    layer_h = Inches(1.15)
    gap = Inches(0.15)
    start_y = Inches(1.3)
    layer_w = CONTENT_WIDTH

    for i, (title, tech, items, color, bg) in enumerate(layers):
        y = start_y + i * (layer_h + gap)
        add_shape(slide, MARGIN_LEFT, y, layer_w, layer_h, fill_color=bg,
                  border_color=color, border_width=Pt(1.5))
        # Layer title
        add_textbox(slide, MARGIN_LEFT + Inches(0.3), y + Inches(0.08),
                    Inches(3.5), Inches(0.35), title, font_size=15,
                    font_color=color, bold=True)
        # Tech stack
        add_textbox(slide, MARGIN_LEFT + Inches(0.3), y + Inches(0.4),
                    Inches(5), Inches(0.3), tech, font_size=11,
                    font_color=TEXT_MUTED)
        # Items (horizontal)
        for j, item in enumerate(items):
            col = j % 5
            x = MARGIN_LEFT + Inches(0.3) + col * Inches(2.3)
            iy = y + Inches(0.68) if j < 5 else y + Inches(0.92)
            add_textbox(slide, x, iy, Inches(2.2), Inches(0.25),
                        f"▸ {item}", font_size=10, font_color=TEXT_PRIMARY)

        # Arrow between layers
        if i < len(layers) - 1:
            arrow_y = y + layer_h + Inches(0.02)
            add_textbox(slide, Inches(6.2), arrow_y, Inches(1), Inches(0.15),
                        "▼", font_size=10, font_color=TEXT_MUTED,
                        alignment=PP_ALIGN.CENTER)


def slide_05_monitoring(prs):
    """실시간 모니터링 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "실시간 모니터링", "Oracle 데이터베이스 성능을 실시간으로 추적하고 분석", PRIMARY_BLUE)
    add_footer(slide, 5)

    features = [
        ("Top SQL 분석", [
            "CPU Time / Elapsed Time 기준 순위",
            "Buffer Gets, Disk Reads 분석",
            "실행 횟수 및 평균 성능 추적",
            "SQL_ID별 상세 드릴다운",
        ], PRIMARY_BLUE, LIGHT_BLUE_BG),
        ("Wait Events", [
            "대기 이벤트 실시간 모니터링",
            "클래스별 분류 및 시각화",
            "Top Wait Events 랭킹",
            "이벤트 발생 추이 그래프",
        ], PRIMARY_CYAN, LIGHT_CYAN_BG),
        ("Session 관리", [
            "활성 세션 실시간 대시보드",
            "세션별 SQL 실행 상태 추적",
            "Blocking Session 감지",
            "세션 Kill 기능",
        ], SUCCESS_GREEN, LIGHT_GREEN_BG),
        ("Lock 모니터링", [
            "Lock Holder/Waiter 트리 시각화",
            "Lock 유형별 분류 및 분석",
            "교착 상태(Deadlock) 감지",
            "Lock 해제 지원",
        ], ORANGE, LIGHT_ORANGE_BG),
    ]

    card_w = Inches(5.65)
    card_h = Inches(2.5)
    gap_x = Inches(0.4)
    gap_y = Inches(0.25)
    start_x = MARGIN_LEFT
    start_y = Inches(1.3)

    for i, (title, items, color, bg) in enumerate(features):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        add_card(slide, x, y, card_w, card_h, title, items,
                 accent_color=color, bg_color=bg, title_size=15, item_size=12)

    # Metrics bar at bottom
    metrics = [
        ("99.9%", "시스템 가용성"),
        ("<100ms", "평균 응답 시간"),
        ("24/7", "실시간 모니터링"),
        ("5초", "데이터 갱신 주기"),
    ]
    mx = MARGIN_LEFT
    my = Inches(6.55)
    mw = Inches(2.7)
    for i, (val, label) in enumerate(metrics):
        add_metric_box(slide, mx + i * (mw + Inches(0.2)), my, mw, Inches(0.8),
                       val, label, PRIMARY_BLUE)


def slide_06_sql_analysis(prs):
    """SQL 고급 분석 & AI 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "SQL 고급 분석 & AI", "AI 기반 스마트 검색, 실행 계획 분석, 패턴 탐지", PRIMARY_PURPLE)
    add_footer(slide, 6)

    tools = [
        ("SQL 통합 검색", "AI 기반 스마트 검색 및 패턴 매칭",
         ["자연어 검색", "정규식 지원", "실시간 필터링", "성능 메트릭 기반 정렬"],
         PRIMARY_BLUE),
        ("실행 계획 분석", "실행 계획 시각화 및 최적화 제안",
         ["트리 시각화", "비용 분석", "인덱스 추천", "Plan Hash 비교"],
         SUCCESS_GREEN),
        ("AI 성능 진단", "머신러닝 기반 성능 이슈 자동 진단",
         ["자동 분석", "개선 제안", "예상 개선율", "LLM 기반 가이드"],
         PRIMARY_PURPLE),
        ("성능 비교 분석", "여러 SQL의 성능 메트릭 비교",
         ["다중 비교", "트렌드 분석", "벤치마킹", "Before/After"],
         ORANGE),
        ("패턴 기반 이슈 탐지", "반복되는 성능 패턴 자동 식별/분류",
         ["패턴 인식", "이상 탐지", "자동 분류", "트렌드 예측"],
         PRIMARY_CYAN),
        ("SQL 클러스터 분석", "유사 SQL 그룹핑 및 A-F 등급",
         ["SQL 유사도 분석", "그룹 최적화", "등급별 관리", "벌크 튜닝"],
         RED_SOFT),
    ]

    cols = 3
    card_w = Inches(3.6)
    card_h = Inches(2.4)
    gap_x = Inches(0.35)
    gap_y = Inches(0.25)
    start_x = (SLIDE_WIDTH - (cols * card_w + (cols - 1) * gap_x)) // 2
    start_y = Inches(1.3)

    for i, (title, desc, items, color) in enumerate(tools):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        add_shape(slide, x, y, card_w, card_h, fill_color=WHITE,
                  border_color=CARD_BORDER, border_width=Pt(1))
        add_rect(slide, x + Inches(0.08), y + Inches(0.08),
                 Inches(0.06), Inches(0.5), color)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.12), card_w - Inches(0.5),
                    Inches(0.3), title, font_size=14, font_color=PRIMARY_DARK, bold=True)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.42), card_w - Inches(0.5),
                    Inches(0.3), desc, font_size=10, font_color=TEXT_MUTED)

        for j, item in enumerate(items):
            add_textbox(slide, x + Inches(0.3), y + Inches(0.8 + j * 0.3),
                        card_w - Inches(0.5), Inches(0.28),
                        f"▸  {item}", font_size=11, font_color=TEXT_PRIMARY)


def slide_07_ai_tuning(prs):
    """AI 튜닝 가이드 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "AI 튜닝 가이드", "온프레미스 LLM 기반 자동 SQL 튜닝 분석 및 가이드", PRIMARY_PURPLE)
    add_footer(slide, 7)

    # Left section - 4 analysis contexts
    contexts = [
        ("SQL 구문 분석", "SQL 텍스트의 구조적 문제점 식별\n조인 방식, 서브쿼리, 힌트 분석"),
        ("실행 계획 분석", "실행 계획 기반 병목 구간 탐지\n비효율적 접근 경로 식별"),
        ("성능 메트릭 분석", "Elapsed Time, CPU Time, Buffer Gets\n실행 횟수 대비 리소스 사용량"),
        ("종합 튜닝 가이드", "분석 결과 통합 리포트\n우선순위별 개선 권고사항"),
    ]

    left_w = Inches(5.8)
    start_y = Inches(1.3)
    card_h = Inches(1.1)
    gap = Inches(0.15)

    for i, (title, desc) in enumerate(contexts):
        y = start_y + i * (card_h + gap)
        add_shape(slide, MARGIN_LEFT, y, left_w, card_h, fill_color=LIGHT_PURPLE_BG,
                  border_color=CARD_BORDER, border_width=Pt(1))
        # Number
        num_shape = add_shape(slide, MARGIN_LEFT + Inches(0.15), y + Inches(0.2),
                              Inches(0.5), Inches(0.5), fill_color=PRIMARY_PURPLE)
        add_textbox(slide, MARGIN_LEFT + Inches(0.15), y + Inches(0.22),
                    Inches(0.5), Inches(0.45), str(i + 1), font_size=18,
                    font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Title & desc
        add_textbox(slide, MARGIN_LEFT + Inches(0.85), y + Inches(0.1),
                    Inches(4.5), Inches(0.3), title, font_size=14,
                    font_color=PRIMARY_DARK, bold=True)
        add_multiline_text(slide, MARGIN_LEFT + Inches(0.85), y + Inches(0.42),
                           Inches(4.5), Inches(0.6), desc.split('\n'),
                           font_size=11, font_color=TEXT_MUTED, line_spacing=1.4)

    # Right section - Key features
    right_x = MARGIN_LEFT + left_w + Inches(0.4)
    right_w = CONTENT_WIDTH - left_w - Inches(0.4)

    add_shape(slide, right_x, start_y, right_w, Inches(4.85),
              fill_color=WHITE, border_color=PRIMARY_PURPLE, border_width=Pt(1.5))
    add_textbox(slide, right_x + Inches(0.3), start_y + Inches(0.15),
                right_w - Inches(0.6), Inches(0.35), "핵심 특징",
                font_size=16, font_color=PRIMARY_PURPLE, bold=True)

    features = [
        "온프레미스 LLM 배포",
        "  - Ollama / vLLM 지원",
        "  - 데이터 외부 유출 방지",
        "",
        "실시간 스트리밍",
        "  - SSE 기반 응답 스트리밍",
        "  - 토큰 단위 실시간 출력",
        "",
        "Clean Architecture",
        "  - Domain / Application /",
        "    Infrastructure 분리",
        "",
        "다중 모델 지원",
        "  - Qwen, Llama, Mistral 등",
        "  - OpenAI 호환 API",
    ]

    fy = start_y + Inches(0.6)
    for line in features:
        is_bold = line and not line.startswith("  ")
        is_empty = line == ""
        if is_empty:
            fy += Inches(0.08)
            continue
        fs = 12 if is_bold else 10
        fc = PRIMARY_DARK if is_bold else TEXT_MUTED
        add_textbox(slide, right_x + Inches(0.3), fy,
                    right_w - Inches(0.6), Inches(0.25),
                    line, font_size=fs, font_color=fc, bold=is_bold)
        fy += Inches(0.22)


def slide_08_execution_plan(prs):
    """실행계획 관리 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "실행계획 관리", "실행 계획 조회, 비교, Baseline 관리, DBMS_XPLAN 통합")
    add_footer(slide, 8)

    features = [
        ("실행계획 조회", [
            "V$SQL_PLAN 기반 실시간 조회",
            "트리 구조 시각화",
            "비용/카디널리티 분석",
            "접근 경로(Access Path) 추적",
            "Predicate 정보 상세 표시",
        ], PRIMARY_BLUE, LIGHT_BLUE_BG),
        ("실행계획 비교", [
            "2개 Plan 동시 비교 뷰",
            "변경된 Operation 하이라이트",
            "성능 메트릭 차이 분석",
            "Plan Hash Value 기반 매칭",
            "Before/After 시각적 비교",
        ], PRIMARY_CYAN, LIGHT_CYAN_BG),
        ("Plan Baseline", [
            "SQL Plan Baseline 등록/관리",
            "Accepted/Fixed Plan 제어",
            "성능 이력 기반 자동 추천",
            "Plan 안정화(Plan Stability)",
            "Baseline 활성화/비활성화",
        ], SUCCESS_GREEN, LIGHT_GREEN_BG),
        ("DBMS_XPLAN", [
            "DBMS_XPLAN.DISPLAY_CURSOR",
            "ALL/TYPICAL/BASIC 포맷",
            "Adaptive Plan 지원",
            "Bind Variable Peeking 분석",
            "Note 섹션 상세 해석",
        ], PRIMARY_PURPLE, LIGHT_PURPLE_BG),
    ]

    card_w = Inches(2.7)
    card_h = Inches(4.8)
    gap = Inches(0.3)
    start_x = (SLIDE_WIDTH - (4 * card_w + 3 * gap)) // 2
    start_y = Inches(1.3)

    for i, (title, items, color, bg) in enumerate(features):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, start_y, card_w, card_h, title, items,
                 accent_color=color, bg_color=bg, title_size=14, item_size=11)


def slide_09_tuning_workflow(prs):
    """튜닝 워크플로우 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "튜닝 워크플로우", "칸반 보드 기반 5단계 체계적 SQL 튜닝 프로세스", PRIMARY_CYAN)
    add_footer(slide, 9)

    # Workflow stages - horizontal
    stages = [
        ("IDENTIFIED", "식별", "Top SQL에서\n자동 식별", PRIMARY_BLUE),
        ("IN_PROGRESS", "진행 중", "담당자 배정\n튜닝 작업 수행", PRIMARY_CYAN),
        ("REVIEW", "리뷰", "튜닝 결과\n검토 및 승인", ORANGE),
        ("COMPLETED", "완료", "Before/After\n성능 비교 확인", SUCCESS_GREEN),
        ("CANCELLED", "취소", "불필요 항목\n사유 기록", TEXT_MUTED),
    ]

    stage_w = Inches(2.1)
    stage_h = Inches(1.7)
    gap = Inches(0.2)
    start_x = (SLIDE_WIDTH - (5 * stage_w + 4 * gap)) // 2
    stage_y = Inches(1.35)

    for i, (code, label, desc, color) in enumerate(stages):
        x = start_x + i * (stage_w + gap)
        add_shape(slide, x, stage_y, stage_w, stage_h, fill_color=WHITE,
                  border_color=color, border_width=Pt(2))
        # Stage number
        add_rect(slide, x, stage_y, stage_w, Inches(0.06), color)
        add_textbox(slide, x, stage_y + Inches(0.15), stage_w, Inches(0.35),
                    label, font_size=15, font_color=color, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), stage_y + Inches(0.5),
                    stage_w - Inches(0.3), Inches(0.25),
                    code, font_size=9, font_color=TEXT_MUTED,
                    alignment=PP_ALIGN.CENTER)
        add_multiline_text(slide, x + Inches(0.15), stage_y + Inches(0.85),
                           stage_w - Inches(0.3), Inches(0.7),
                           desc.split('\n'), font_size=10, font_color=TEXT_PRIMARY,
                           alignment=PP_ALIGN.CENTER, line_spacing=1.4)
        # Arrow
        if i < 4:
            arrow_x = x + stage_w + Inches(0.02)
            add_textbox(slide, arrow_x, stage_y + Inches(0.6), Inches(0.18), Inches(0.3),
                        "→", font_size=16, font_color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Bottom section - key features
    bottom_y = Inches(3.4)
    cards = [
        ("칸반 보드", [
            "드래그 & 드롭 상태 변경",
            "우선순위별 정렬 (Critical/High/Medium/Low)",
            "담당자별 필터링",
            "실시간 업데이트",
        ], PRIMARY_BLUE, LIGHT_BLUE_BG),
        ("Before/After 비교", [
            "Elapsed Time 개선율",
            "CPU Time 개선율",
            "Buffer Gets 개선율",
            "Plan Hash Value 변경 추적",
        ], SUCCESS_GREEN, LIGHT_GREEN_BG),
        ("협업 기능", [
            "튜닝 태스크 댓글/토론",
            "멘션(@) 기능",
            "첨부 파일 지원",
            "작업 이력 타임라인",
        ], PRIMARY_PURPLE, LIGHT_PURPLE_BG),
    ]

    card_w = Inches(3.6)
    card_h = Inches(2.6)
    gap = Inches(0.35)
    start_x = (SLIDE_WIDTH - (3 * card_w + 2 * gap)) // 2

    for i, (title, items, color, bg) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, bottom_y, card_w, card_h, title, items,
                 accent_color=color, bg_color=bg, title_size=14, item_size=11)


def slide_10_performance_history(prs):
    """성능 히스토리 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "성능 히스토리", "A-F 등급 시스템 기반 일별/시간대별 SQL 성능 추이 분석", SUCCESS_GREEN)
    add_footer(slide, 10)

    # Grade system
    grades = [
        ("A", "최적", "0-50ms", RGBColor(0x22, 0xC5, 0x5E)),
        ("B", "양호", "50-200ms", RGBColor(0x84, 0xCC, 0x16)),
        ("C", "보통", "200-500ms", RGBColor(0xEA, 0xB3, 0x08)),
        ("D", "주의", "500ms-1s", RGBColor(0xF9, 0x73, 0x16)),
        ("E", "경고", "1s-3s", RGBColor(0xEF, 0x44, 0x44)),
        ("F", "위험", "3s+", RGBColor(0xDC, 0x26, 0x26)),
    ]

    grade_w = Inches(1.75)
    grade_h = Inches(1.2)
    gap = Inches(0.2)
    start_x = (SLIDE_WIDTH - (6 * grade_w + 5 * gap)) // 2
    start_y = Inches(1.3)

    for i, (grade, label, threshold, color) in enumerate(grades):
        x = start_x + i * (grade_w + gap)
        add_shape(slide, x, start_y, grade_w, grade_h, fill_color=WHITE,
                  border_color=color, border_width=Pt(2))
        add_textbox(slide, x, start_y + Inches(0.1), grade_w, Inches(0.45),
                    grade, font_size=30, font_color=color, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, start_y + Inches(0.55), grade_w, Inches(0.25),
                    label, font_size=12, font_color=PRIMARY_DARK,
                    alignment=PP_ALIGN.CENTER, bold=True)
        add_textbox(slide, x, start_y + Inches(0.8), grade_w, Inches(0.25),
                    threshold, font_size=10, font_color=TEXT_MUTED,
                    alignment=PP_ALIGN.CENTER)

    # Features below
    bottom_features = [
        ("일별 성능 추이", [
            "SQL별 일별 성능 등급 변화",
            "등급 분포 히스토그램",
            "개선/악화 SQL 자동 식별",
            "기간별 비교 분석",
        ], PRIMARY_BLUE, LIGHT_BLUE_BG),
        ("시간대별 분석", [
            "24시간 히트맵 시각화",
            "피크 시간대 자동 탐지",
            "시간대별 Top SQL 변화",
            "업무 시간대 집중 분석",
        ], PRIMARY_CYAN, LIGHT_CYAN_BG),
        ("수집 설정", [
            "수집 주기 커스터마이징 (5s~1h)",
            "대상 SQL 필터 설정",
            "데이터 보존 기간 관리",
            "자동 정리(Purge) 정책",
        ], PRIMARY_PURPLE, LIGHT_PURPLE_BG),
    ]

    card_w = Inches(3.6)
    card_h = Inches(2.5)
    gap = Inches(0.35)
    start_x = (SLIDE_WIDTH - (3 * card_w + 2 * gap)) // 2
    by = Inches(2.85)

    for i, (title, items, color, bg) in enumerate(bottom_features):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, by, card_w, card_h, title, items,
                 accent_color=color, bg_color=bg, title_size=14, item_size=11)


def slide_11_awr_addm(prs):
    """AWR/ADDM/STATSPACK 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "AWR / ADDM / STATSPACK", "Oracle Enterprise Edition & Standard Edition 모두 지원", ORANGE)
    add_footer(slide, 11)

    # Enterprise Edition section
    add_shape(slide, MARGIN_LEFT, Inches(1.3), Inches(5.6), Inches(5.4),
              fill_color=LIGHT_BLUE_BG, border_color=PRIMARY_BLUE, border_width=Pt(1.5))
    add_rect(slide, MARGIN_LEFT, Inches(1.3), Inches(5.6), Inches(0.5), PRIMARY_BLUE)
    add_textbox(slide, MARGIN_LEFT + Inches(0.3), Inches(1.35), Inches(5), Inches(0.4),
                "Enterprise Edition", font_size=16, font_color=WHITE, bold=True)

    ee_features = [
        ("AWR (Automatic Workload Repository)", [
            "자동 스냅샷 수집 및 관리",
            "AWR 리포트 생성 및 비교",
            "기간별 워크로드 분석",
            "Top Events / SQL / Segments",
        ]),
        ("ADDM (Automatic Database Diagnostic Monitor)", [
            "자동 성능 진단 보고서",
            "Finding / Recommendation 분석",
            "영향도(Impact) 기반 우선순위",
            "이전 기간 대비 변화 분석",
        ]),
        ("ASH (Active Session History)", [
            "실시간 세션 히스토리 분석",
            "시간대별 세션 분포",
            "SQL/Event/Module별 분류",
            "ASH 리포트 생성",
        ]),
    ]

    ey = Inches(2.0)
    for title, items in ee_features:
        add_textbox(slide, MARGIN_LEFT + Inches(0.3), ey, Inches(5), Inches(0.3),
                    title, font_size=13, font_color=PRIMARY_DARK, bold=True)
        ey += Inches(0.32)
        for item in items:
            add_textbox(slide, MARGIN_LEFT + Inches(0.5), ey, Inches(4.8), Inches(0.25),
                        f"▸  {item}", font_size=10, font_color=TEXT_PRIMARY)
            ey += Inches(0.24)
        ey += Inches(0.15)

    # Standard Edition section
    se_x = MARGIN_LEFT + Inches(6.0)
    se_w = CONTENT_WIDTH - Inches(6.0)
    add_shape(slide, se_x, Inches(1.3), se_w, Inches(5.4),
              fill_color=LIGHT_GREEN_BG, border_color=SUCCESS_GREEN, border_width=Pt(1.5))
    add_rect(slide, se_x, Inches(1.3), se_w, Inches(0.5), SUCCESS_GREEN)
    add_textbox(slide, se_x + Inches(0.3), Inches(1.35), se_w - Inches(0.6), Inches(0.4),
                "Standard Edition", font_size=16, font_color=WHITE, bold=True)

    se_features = [
        ("STATSPACK", [
            "AWR 대안으로 무료 사용",
            "스냅샷 수집/삭제 관리",
            "STATSPACK 리포트 생성",
            "Level 설정 (1-10)",
            "기간별 비교 분석",
        ]),
        ("통계정보 수집", [
            "테이블/인덱스 통계 수집",
            "수집 스케줄 관리",
            "Stale 통계 자동 감지",
            "수집 이력 조회",
        ]),
        ("SQL Trace", [
            "세션별 SQL Trace 활성화",
            "TKPROF 분석 지원",
            "10046/10053 이벤트",
            "Trace 파일 관리",
        ]),
    ]

    sy = Inches(2.0)
    for title, items in se_features:
        add_textbox(slide, se_x + Inches(0.3), sy, se_w - Inches(0.6), Inches(0.3),
                    title, font_size=13, font_color=PRIMARY_DARK, bold=True)
        sy += Inches(0.32)
        for item in items:
            add_textbox(slide, se_x + Inches(0.5), sy, se_w - Inches(0.8), Inches(0.25),
                        f"▸  {item}", font_size=10, font_color=TEXT_PRIMARY)
            sy += Inches(0.24)
        sy += Inches(0.15)


def slide_12_reports(prs):
    """리포트 시스템 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "리포트 시스템", "템플릿 기반 리포트 생성, 스케줄링, 다중 포맷 내보내기")
    add_footer(slide, 12)

    sections = [
        ("리포트 템플릿", [
            "일별/주별/월별 성능 리포트",
            "튜닝 현황 요약 리포트",
            "SQL 성능 변화 리포트",
            "AWR 기반 워크로드 리포트",
            "커스텀 템플릿 생성",
        ], PRIMARY_BLUE, LIGHT_BLUE_BG),
        ("스케줄링", [
            "Cron 기반 자동 생성",
            "이메일 자동 발송",
            "수신자 그룹 관리",
            "생성 이력 조회",
            "실행 시간 모니터링",
        ], PRIMARY_PURPLE, LIGHT_PURPLE_BG),
        ("내보내기", [
            "Excel (.xlsx) 내보내기",
            "PDF 리포트 생성",
            "HTML 웹 리포트",
            "차트/그래프 포함",
            "다국어 지원 (한/영)",
        ], SUCCESS_GREEN, LIGHT_GREEN_BG),
    ]

    card_w = Inches(3.6)
    card_h = Inches(3.2)
    gap = Inches(0.35)
    start_x = (SLIDE_WIDTH - (3 * card_w + 2 * gap)) // 2
    start_y = Inches(1.5)

    for i, (title, items, color, bg) in enumerate(sections):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, start_y, card_w, card_h, title, items,
                 accent_color=color, bg_color=bg, title_size=15, item_size=12)

    # Bottom highlight
    add_shape(slide, MARGIN_LEFT, Inches(5.2), CONTENT_WIDTH, Inches(1.4),
              fill_color=WHITE, border_color=CARD_BORDER, border_width=Pt(1))
    add_textbox(slide, MARGIN_LEFT + Inches(0.4), Inches(5.35), Inches(3), Inches(0.35),
                "리포트 워크플로우", font_size=15, font_color=PRIMARY_DARK, bold=True)

    steps = ["템플릿 선택", "기간/대상 설정", "리포트 생성", "검토/승인", "배포/공유"]
    step_w = Inches(2.0)
    sx = MARGIN_LEFT + Inches(0.3)
    sy = Inches(5.85)
    for i, step in enumerate(steps):
        add_shape(slide, sx + i * (step_w + Inches(0.15)), sy,
                  step_w, Inches(0.5), fill_color=LIGHT_BLUE_BG,
                  border_color=PRIMARY_BLUE, border_width=Pt(1))
        add_textbox(slide, sx + i * (step_w + Inches(0.15)), sy + Inches(0.08),
                    step_w, Inches(0.35), f"{i+1}. {step}", font_size=11,
                    font_color=PRIMARY_DARK, alignment=PP_ALIGN.CENTER, bold=True)
        if i < 4:
            add_textbox(slide, sx + i * (step_w + Inches(0.15)) + step_w,
                        sy + Inches(0.05), Inches(0.15), Inches(0.35),
                        "→", font_size=14, font_color=TEXT_MUTED,
                        alignment=PP_ALIGN.CENTER)


def slide_13_db_connections(prs):
    """DB 연결 관리 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "DB 연결 관리", "다중 Oracle 데이터베이스 연결, AES-256 암호화, Health Check", SUCCESS_GREEN)
    add_footer(slide, 13)

    features = [
        ("다중 연결 관리", [
            "여러 Oracle DB 동시 연결",
            "연결별 별칭(Alias) 관리",
            "TNS / Easy Connect 지원",
            "연결 풀링 최적화",
            "실시간 연결 상태 표시",
        ], PRIMARY_BLUE, LIGHT_BLUE_BG),
        ("보안 암호화", [
            "AES-256-CBC 비밀번호 암호화",
            "환경변수 기반 키 관리",
            "비밀번호 마스킹 표시",
            "암호화 키 로테이션",
            "Zero-Knowledge 설계",
        ], SUCCESS_GREEN, LIGHT_GREEN_BG),
        ("Health Check", [
            "주기적 연결 상태 점검",
            "Oracle Edition 자동 감지",
            "Enterprise/Standard 구분",
            "연결 실패 자동 알림",
            "연결 복구 자동 재시도",
        ], PRIMARY_CYAN, LIGHT_CYAN_BG),
    ]

    card_w = Inches(3.6)
    card_h = Inches(3.5)
    gap = Inches(0.35)
    start_x = (SLIDE_WIDTH - (3 * card_w + 2 * gap)) // 2
    start_y = Inches(1.4)

    for i, (title, items, color, bg) in enumerate(features):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, start_y, card_w, card_h, title, items,
                 accent_color=color, bg_color=bg, title_size=15, item_size=12)

    # Connection info schema
    add_shape(slide, MARGIN_LEFT, Inches(5.3), CONTENT_WIDTH, Inches(1.3),
              fill_color=WHITE, border_color=CARD_BORDER, border_width=Pt(1))
    add_textbox(slide, MARGIN_LEFT + Inches(0.4), Inches(5.4), Inches(4), Inches(0.3),
                "연결 정보 스키마", font_size=14, font_color=PRIMARY_DARK, bold=True)
    schema_items = [
        "호스트(Host) / 포트(Port) / SID or Service Name",
        "사용자명(Username) / 암호화된 비밀번호",
        "Oracle Edition 정보 / 연결 상태 / 마지막 점검 시간",
    ]
    for i, item in enumerate(schema_items):
        add_textbox(slide, MARGIN_LEFT + Inches(0.4), Inches(5.75 + i * 0.25),
                    CONTENT_WIDTH - Inches(0.8), Inches(0.25),
                    f"▸  {item}", font_size=11, font_color=TEXT_PRIMARY)


def slide_14_ai_detail(prs):
    """AI/LLM 통합 상세 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "AI/LLM 통합 상세", "6가지 Use Case, Clean Architecture, 온프레미스 보안", PRIMARY_PURPLE)
    add_footer(slide, 14)

    # 6 Use Cases
    use_cases = [
        ("SQL 구조 분석", "SQL 텍스트의 구조적\n문제점 자동 식별", PRIMARY_BLUE),
        ("실행계획 해석", "실행계획 기반\n병목 구간 분석", PRIMARY_CYAN),
        ("튜닝 가이드", "단계별 최적화\n가이드 생성", PRIMARY_PURPLE),
        ("인덱스 추천", "누락/비효율 인덱스\n자동 추천", SUCCESS_GREEN),
        ("쿼리 리팩토링", "AI 기반 SQL\n재작성 제안", ORANGE),
        ("패턴 분석", "성능 이슈 패턴\n자동 분류/학습", RED_SOFT),
    ]

    uc_w = Inches(1.8)
    uc_h = Inches(1.4)
    gap = Inches(0.2)
    start_x = (SLIDE_WIDTH - (6 * uc_w + 5 * gap)) // 2
    start_y = Inches(1.3)

    for i, (title, desc, color) in enumerate(use_cases):
        x = start_x + i * (uc_w + gap)
        add_shape(slide, x, start_y, uc_w, uc_h, fill_color=WHITE,
                  border_color=color, border_width=Pt(1.5))
        add_rect(slide, x, start_y, uc_w, Inches(0.05), color)
        add_textbox(slide, x + Inches(0.1), start_y + Inches(0.12),
                    uc_w - Inches(0.2), Inches(0.3), title, font_size=11,
                    font_color=PRIMARY_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        add_multiline_text(slide, x + Inches(0.1), start_y + Inches(0.45),
                           uc_w - Inches(0.2), Inches(0.8), desc.split('\n'),
                           font_size=9, font_color=TEXT_MUTED, alignment=PP_ALIGN.CENTER,
                           line_spacing=1.3)

    # Clean Architecture & On-premises
    bottom_y = Inches(3.0)

    # Clean Architecture
    add_shape(slide, MARGIN_LEFT, bottom_y, Inches(5.6), Inches(3.5),
              fill_color=WHITE, border_color=PRIMARY_PURPLE, border_width=Pt(1.5))
    add_textbox(slide, MARGIN_LEFT + Inches(0.3), bottom_y + Inches(0.15),
                Inches(5), Inches(0.35), "Clean Architecture 설계",
                font_size=15, font_color=PRIMARY_PURPLE, bold=True)

    layers = [
        ("Domain Layer", "엔티티, 값 객체, 도메인 서비스", LIGHT_PURPLE_BG),
        ("Application Layer", "Use Case, DTO, 포트 정의", LIGHT_BLUE_BG),
        ("Infrastructure Layer", "LLM 클라이언트, 리포지토리 구현", LIGHT_GREEN_BG),
        ("Presentation Layer", "React 컴포넌트, API Route 핸들러", LIGHT_CYAN_BG),
    ]

    ly = bottom_y + Inches(0.6)
    for title, desc, bg in layers:
        add_shape(slide, MARGIN_LEFT + Inches(0.3), ly, Inches(4.9), Inches(0.6),
                  fill_color=bg, border_color=CARD_BORDER, border_width=Pt(1))
        add_textbox(slide, MARGIN_LEFT + Inches(0.5), ly + Inches(0.05),
                    Inches(2), Inches(0.25), title, font_size=11,
                    font_color=PRIMARY_DARK, bold=True)
        add_textbox(slide, MARGIN_LEFT + Inches(0.5), ly + Inches(0.3),
                    Inches(4.3), Inches(0.25), desc, font_size=10,
                    font_color=TEXT_MUTED)
        ly += Inches(0.68)

    # On-premises Security
    sec_x = MARGIN_LEFT + Inches(6.0)
    sec_w = CONTENT_WIDTH - Inches(6.0)
    add_shape(slide, sec_x, bottom_y, sec_w, Inches(3.5),
              fill_color=WHITE, border_color=SUCCESS_GREEN, border_width=Pt(1.5))
    add_textbox(slide, sec_x + Inches(0.3), bottom_y + Inches(0.15),
                sec_w - Inches(0.6), Inches(0.35), "온프레미스 보안",
                font_size=15, font_color=SUCCESS_GREEN, bold=True)

    security_items = [
        "데이터가 외부로 전송되지 않음",
        "자체 서버에서 LLM 실행",
        "Ollama / vLLM 엔진 지원",
        "OpenAI 호환 API 인터페이스",
        "모델 선택 자유도 (Qwen, Llama, ...)",
        "스트리밍 응답 (SSE)",
        "대화 이력 관리",
        "응답 내보내기 기능",
    ]
    sy = bottom_y + Inches(0.65)
    for item in security_items:
        add_textbox(slide, sec_x + Inches(0.3), sy, sec_w - Inches(0.6), Inches(0.28),
                    f"▸  {item}", font_size=11, font_color=TEXT_PRIMARY)
        sy += Inches(0.32)


def slide_15_security(prs):
    """보안 & 접근 제어 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "보안 & 접근 제어", "JWT 인증, bcrypt 해싱, RBAC, 감사 로그", SUCCESS_GREEN)
    add_footer(slide, 15)

    sections = [
        ("인증 (Authentication)", [
            "NextAuth v4 JWT 전략",
            "bcrypt 비밀번호 해싱",
            "세션 토큰 자동 갱신",
            "로그인 시도 제한",
            "비밀번호 정책 적용",
        ], PRIMARY_BLUE, LIGHT_BLUE_BG),
        ("인가 (Authorization)", [
            "Role 기반 접근 제어 (RBAC)",
            "Admin / DBA / Viewer 역할",
            "페이지별 권한 설정",
            "API 엔드포인트 보호",
            "동적 메뉴 권한 적용",
        ], PRIMARY_PURPLE, LIGHT_PURPLE_BG),
        ("데이터 보안", [
            "AES-256-CBC DB 비밀번호 암호화",
            "환경변수 기반 시크릿 관리",
            "HTTPS 전송 암호화",
            "SQL Injection 방지 (Drizzle ORM)",
            "XSS/CSRF 보호",
        ], SUCCESS_GREEN, LIGHT_GREEN_BG),
        ("감사 & 모니터링", [
            "전체 작업 이력 추적",
            "로그인/로그아웃 기록",
            "DB 연결 변경 이력",
            "튜닝 작업 감사 로그",
            "타임스탬프 기반 추적",
        ], ORANGE, LIGHT_ORANGE_BG),
    ]

    card_w = Inches(2.7)
    card_h = Inches(3.5)
    gap = Inches(0.3)
    start_x = (SLIDE_WIDTH - (4 * card_w + 3 * gap)) // 2
    start_y = Inches(1.4)

    for i, (title, items, color, bg) in enumerate(sections):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, start_y, card_w, card_h, title, items,
                 accent_color=color, bg_color=bg, title_size=13, item_size=11)

    # Trust indicators
    indicators = [
        ("ISO 27001", "보안 인증 기준 준수"),
        ("99.9%", "시스템 가용성"),
        ("AES-256", "군사급 암호화"),
        ("RBAC", "역할 기반 접근 제어"),
    ]

    iy = Inches(5.3)
    iw = Inches(2.7)
    gap = Inches(0.3)
    ix = (SLIDE_WIDTH - (4 * iw + 3 * gap)) // 2

    for i, (val, label) in enumerate(indicators):
        x = ix + i * (iw + gap)
        add_metric_box(slide, x, iy, iw, Inches(0.85), val, label, SUCCESS_GREEN)


def slide_16_competitive(prs):
    """경쟁 우위 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "경쟁 우위", "기존 도구 대비 Narae TMS v2.0의 차별화 포인트", PRIMARY_BLUE)
    add_footer(slide, 16)

    # Comparison table header
    table_x = MARGIN_LEFT
    table_y = Inches(1.4)
    col_w = [Inches(3.0), Inches(3.5), Inches(5.2)]
    row_h = Inches(0.55)
    headers = ["기능 영역", "기존 도구", "Narae TMS v2.0"]
    header_colors = [TEXT_MUTED, RED_SOFT, PRIMARY_BLUE]

    # Header row
    cx = table_x
    for i, (header, color) in enumerate(zip(headers, header_colors)):
        bg = PRIMARY_DARK if i == 0 else WHITE
        fc = WHITE if i == 0 else color
        add_shape(slide, cx, table_y, col_w[i], row_h,
                  fill_color=bg, border_color=CARD_BORDER, border_width=Pt(1))
        add_textbox(slide, cx + Inches(0.15), table_y + Inches(0.1),
                    col_w[i] - Inches(0.3), Inches(0.35),
                    header, font_size=13, font_color=fc, bold=True,
                    alignment=PP_ALIGN.CENTER)
        cx += col_w[i]

    # Data rows
    rows = [
        ("모니터링", "수동 SQL*Plus / EM", "실시간 대시보드 + AI 자동 알림"),
        ("SQL 튜닝", "개인 노하우 의존", "AI 가이드 + 체계적 워크플로우"),
        ("이력 관리", "엑셀/문서 수동 관리", "자동 Before/After + 타임라인"),
        ("실행계획", "DBMS_XPLAN 수동 실행", "시각화 + 비교 + Baseline 관리"),
        ("보안", "공유 계정 / 평문 저장", "RBAC + AES-256 + 감사 로그"),
        ("협업", "이메일/메신저 비체계적", "칸반 보드 + 댓글 + 멘션"),
        ("리포트", "수동 생성 / 비정형", "자동 템플릿 + 스케줄 + 다중 포맷"),
        ("AI/LLM", "미지원", "온프레미스 LLM + 6가지 Use Case"),
        ("Edition 지원", "EE 중심", "Enterprise + Standard 모두 지원"),
    ]

    for i, (feature, old_tool, tms) in enumerate(rows):
        ry = table_y + (i + 1) * row_h
        bg = WHITE if i % 2 == 0 else BG_LIGHT
        cx = table_x
        for j, text in enumerate([feature, old_tool, tms]):
            add_shape(slide, cx, ry, col_w[j], row_h,
                      fill_color=bg, border_color=CARD_BORDER, border_width=Pt(0.5))
            fc = PRIMARY_DARK if j == 0 else (TEXT_MUTED if j == 1 else PRIMARY_BLUE)
            b = j == 0
            add_textbox(slide, cx + Inches(0.15), ry + Inches(0.1),
                        col_w[j] - Inches(0.3), Inches(0.35),
                        text, font_size=11, font_color=fc, bold=b)
            cx += col_w[j]


def slide_17_tech_stack(prs):
    """기술 스택 슬라이드"""
    slide = make_slide(prs)
    add_title_bar(slide, "기술 스택", "모던 웹 기술 기반 엔터프라이즈급 아키텍처")
    add_footer(slide, 17)

    stacks = [
        ("Frontend", PRIMARY_BLUE, LIGHT_BLUE_BG, [
            "Next.js 15.1 (App Router)",
            "React 19",
            "TypeScript (Strict)",
            "Tailwind CSS",
            "Shadcn UI + Radix UI",
            "React Query",
            "Zustand",
            "React Hook Form + Zod",
            "Framer Motion",
            "Lucide Icons",
        ]),
        ("Backend", PRIMARY_CYAN, LIGHT_CYAN_BG, [
            "Next.js API Routes",
            "NextAuth v4 (JWT)",
            "Drizzle ORM",
            "PostgreSQL 17",
            "Oracle Database (대상)",
            "bcrypt (해싱)",
            "AES-256-CBC (암호화)",
            "node-oracledb",
            "ts-pattern",
            "es-toolkit",
        ]),
        ("AI / LLM", PRIMARY_PURPLE, LIGHT_PURPLE_BG, [
            "Ollama (온프레미스)",
            "vLLM (대안 엔진)",
            "OpenAI 호환 API",
            "SSE 스트리밍",
            "Clean Architecture",
            "도메인 주도 설계",
            "Qwen / Llama / Mistral",
            "프롬프트 엔지니어링",
            "대화 이력 관리",
            "응답 내보내기",
        ]),
        ("DevOps & Export", SUCCESS_GREEN, LIGHT_GREEN_BG, [
            "Turbopack (개발)",
            "ESLint / Prettier",
            "Drizzle Migrations",
            "GitHub CI/CD",
            "Excel 내보내기 (xlsx)",
            "PDF 생성",
            "HTML 리포트",
            "next-themes (다크모드)",
            "date-fns",
            "react-use",
        ]),
    ]

    card_w = Inches(2.7)
    card_h = Inches(5.2)
    gap = Inches(0.3)
    start_x = (SLIDE_WIDTH - (4 * card_w + 3 * gap)) // 2
    start_y = Inches(1.3)

    for i, (title, color, bg, items) in enumerate(stacks):
        x = start_x + i * (card_w + gap)
        add_shape(slide, x, start_y, card_w, card_h, fill_color=bg,
                  border_color=color, border_width=Pt(1.5))
        # Title bar
        add_rect(slide, x, start_y, card_w, Inches(0.5), color)
        add_textbox(slide, x, start_y + Inches(0.08), card_w, Inches(0.35),
                    title, font_size=15, font_color=WHITE, bold=True,
                    alignment=PP_ALIGN.CENTER)
        # Items
        iy = start_y + Inches(0.65)
        for item in items:
            add_textbox(slide, x + Inches(0.2), iy, card_w - Inches(0.4), Inches(0.28),
                        f"▸  {item}", font_size=11, font_color=TEXT_PRIMARY)
            iy += Inches(0.4)


def slide_18_closing(prs):
    """감사합니다 슬라이드"""
    slide = make_slide(prs)
    set_slide_bg(slide, PRIMARY_DARK)

    # Decorative elements
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PRIMARY_BLUE)
    add_shape(slide, Inches(-2), Inches(5), Inches(6), Inches(6),
              fill_color=RGBColor(0x1E, 0x29, 0x3B))

    # Thank you text
    add_textbox(slide, Inches(0), Inches(1.5), SLIDE_WIDTH, Inches(0.8),
                "감사합니다", font_size=48, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0), Inches(2.4), SLIDE_WIDTH, Inches(0.5),
                "Thank You", font_size=24, font_color=PRIMARY_CYAN,
                alignment=PP_ALIGN.CENTER)

    # Separator
    add_rect(slide, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.04), PRIMARY_BLUE)

    # Product info
    add_textbox(slide, Inches(0), Inches(3.6), SLIDE_WIDTH, Inches(0.5),
                PRODUCT_NAME, font_size=22, font_color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0), Inches(4.1), SLIDE_WIDTH, Inches(0.4),
                "Enterprise-Grade Oracle SQL Tuning Management System",
                font_size=14, font_color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    # Contact info
    contact_items = [
        f"{COMPANY_NAME}",
        "www.narae-it.co.kr",
        "contact@narae-it.co.kr",
    ]
    cy = Inches(5.0)
    for item in contact_items:
        add_textbox(slide, Inches(0), cy, SLIDE_WIDTH, Inches(0.35),
                    item, font_size=14, font_color=RGBColor(0x94, 0xA3, 0xB8),
                    alignment=PP_ALIGN.CENTER)
        cy += Inches(0.35)

    # Bottom tagline
    add_textbox(slide, Inches(0), Inches(6.4), SLIDE_WIDTH, Inches(0.3),
                "실시간 모니터링  •  AI 튜닝 가이드  •  체계적 워크플로우  •  엔터프라이즈 보안",
                font_size=11, font_color=PRIMARY_BLUE, alignment=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────

def main():
    prs = Presentation()

    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Generating Narae TMS v2.0 Introduction PPT...")

    slide_builders = [
        ("표지", slide_01_cover),
        ("문제 정의", slide_02_problem),
        ("솔루션 개요", slide_03_solution),
        ("시스템 아키텍처", slide_04_architecture),
        ("실시간 모니터링", slide_05_monitoring),
        ("SQL 고급 분석 & AI", slide_06_sql_analysis),
        ("AI 튜닝 가이드", slide_07_ai_tuning),
        ("실행계획 관리", slide_08_execution_plan),
        ("튜닝 워크플로우", slide_09_tuning_workflow),
        ("성능 히스토리", slide_10_performance_history),
        ("AWR/ADDM/STATSPACK", slide_11_awr_addm),
        ("리포트 시스템", slide_12_reports),
        ("DB 연결 관리", slide_13_db_connections),
        ("AI/LLM 통합 상세", slide_14_ai_detail),
        ("보안 & 접근 제어", slide_15_security),
        ("경쟁 우위", slide_16_competitive),
        ("기술 스택", slide_17_tech_stack),
        ("감사합니다", slide_18_closing),
    ]

    for i, (name, builder) in enumerate(slide_builders, 1):
        print(f"  [{i:2d}/18] {name}...")
        builder(prs)

    output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)),
                               "narae-tms-v2-introduction.pptx")
    prs.save(output_path)
    print(f"\nDone! Output: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
